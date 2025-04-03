import streamlit as st
import pandas as pd
import random
from faker import Faker
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

fake = Faker("en_CA")

products = [
    {"item": "Milk", "price": 3.49}, {"item": "Bread", "price": 2.99},
    {"item": "Eggs", "price": 4.99}, {"item": "Apples", "price": 1.29},
    {"item": "Bananas", "price": 0.79}, {"item": "Chicken Breast", "price": 6.49},
    {"item": "Toilet Paper", "price": 5.99}, {"item": "Shampoo", "price": 4.50},
    {"item": "Orange Juice", "price": 3.99}, {"item": "Chips", "price": 2.49},
]

regions = ["Ontario", "Quebec", "British Columbia", "Alberta"]
incomes = ["<40K", "40K–70K", "70K–100K", "100K+"]

def generate_respondents(n, region):
    return pd.DataFrame([{
        "respondent_id": f"{region}_{i}",
        "region": region,
        "age": random.randint(18, 80),
        "income_bracket": random.choices(incomes, weights=[0.2, 0.4, 0.3, 0.1])[0],
        "loyalty": random.choice([True, False])
    } for i in range(n)])

def generate_basket(row, promo_dict):
    basket = []
    for item in random.sample(products, random.randint(3, 6)):
        qty = random.randint(1, 3)
        price = item["price"]
        promo = promo_dict.get(item["item"], 0) / 100
        if row["loyalty"]: price *= 0.9
        price *= (1 - promo)
        basket.append({
            "respondent_id": row["respondent_id"], "region": row["region"],
            "age": row["age"], "income_bracket": row["income_bracket"],
            "loyalty": row["loyalty"], "item": item["item"],
            "quantity": qty, "unit_price": round(price, 2),
            "total_price": round(qty * price, 2)
        })
    return basket

def generate_insights(before, after, promo_dict, region):
    insights = [f"📍 **{region} Promo Summary**"]
    for item, promo in promo_dict.items():
        b = before[(before["item"] == item) & (before["region"] == region)]["total_price"].sum()
        a = after[(after["item"] == item) & (after["region"] == region)]["total_price"].sum()
        if b > 0:
            change = ((a - b) / b) * 100
            insights.append(f"🛍️ *{item}* sales increased by **{change:.1f}%** under a {promo}% promo.")
    reg_total = after[after["region"] == region]["total_price"].sum()
    insights.append(f"💵 *{region}* total spend was **${reg_total:.2f}**.")
    loyal_avg = after[(after["region"] == region) & (after["loyalty"] == True)]["total_price"].mean()
    nonloyal_avg = after[(after["region"] == region) & (after["loyalty"] == False)]["total_price"].mean()
    if not pd.isna(loyal_avg) and not pd.isna(nonloyal_avg):
        diff = loyal_avg - nonloyal_avg
        insights.append(f"🎯 Loyalty members spent **${diff:.2f}** more on average.")
    return insights

def apply_brand_style(slide, title_text, bullet_points):
    slide.shapes.title.text = title_text
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    for bullet in bullet_points:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(50, 50, 50)

def create_multi_region_ppt(insight_blocks):
    ppt = Presentation()
    title = ppt.slides.add_slide(ppt.slide_layouts[0])
    title.shapes.title.text = "Multi-Market Campaign Report"
    title.placeholders[1].text = "Regional Performance Breakdown"
    for block in insight_blocks:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        label = block[0].replace("📍", "").strip()
        apply_brand_style(slide, label, block[1:])
    output = io.BytesIO()
    ppt.save(output)
    return output.getvalue()

st.title("🌐 Multi-Market Campaign Simulator")

num = st.slider("Respondents per Region", 100, 1000, 300)
selected_regions = st.multiselect("Regions to Simulate", regions, default=regions[:2])

promo_items = st.multiselect("Promo Items", [p["item"] for p in products])
region_promos = {}
for region in selected_regions:
    with st.expander(f"{region} Promos"):
        region_promos[region] = {
            item: st.slider(f"{item} Discount in {region}", 0, 50, 0, key=f"{region}_{item}")
            for item in promo_items
        }

if st.button("Run Regional Campaigns"):
    all_before = []
    all_after = []
    all_insights = []

    for region in selected_regions:
        df_r = generate_respondents(num, region)
        before = pd.DataFrame([row for _, r in df_r.iterrows() for row in generate_basket(r, {})])
        after = pd.DataFrame([row for _, r in df_r.iterrows() for row in generate_basket(r, region_promos[region])])
        all_before.append(before)
        all_after.append(after)
        all_insights.append(generate_insights(before, after, region_promos[region], region))

    df_before = pd.concat(all_before)
    df_after = pd.concat(all_after)

    for insights in all_insights:
        st.subheader(insights[0])
        for line in insights[1:]:
            st.markdown(line)

    ppt = create_multi_region_ppt(all_insights)
    st.download_button("📥 Download Multi-Market Report", ppt, file_name="regional_campaign_report.pptx")
